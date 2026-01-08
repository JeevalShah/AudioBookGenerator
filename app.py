import streamlit as st
import os
from dotenv import load_dotenv
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from num2words import num2words
import re
import google.generativeai as genai
from gtts import gTTS

# Metrics libraries
import speech_recognition as sr
from pydub import AudioSegment
import jiwer
import textstat
import string

# Supported languages
SUPPORTED_LANGUAGES = {
    "English": {"code": "en", "stt_code": "en-US"},
    "Hindi": {"code": "hi", "stt_code": "hi-IN"},
    "Spanish": {"code": "es", "stt_code": "es-ES"},
    "French": {"code": "fr", "stt_code": "fr-FR"},
    "German": {"code": "de", "stt_code": "de-DE"},
    "Italian": {"code": "it", "stt_code": "it-IT"},
    "Portuguese": {"code": "pt", "stt_code": "pt-PT"},
    "Russian": {"code": "ru", "stt_code": "ru-RU"},
    "Arabic": {"code": "ar", "stt_code": "ar-SA"},
    "Chinese (Simplified)": {"code": "zh", "stt_code": "zh-CN"},
    "Japanese": {"code": "ja", "stt_code": "ja-JP"},
    "Korean": {"code": "ko", "stt_code": "ko-KR"}
}

# Load API keys
load_dotenv(".env")
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")

if not GEMINI_API_KEY:
    st.error("GEMINI API key not found. Add GEMINI_API_KEY to .env")
    st.stop()

# Configure Gemini 
genai.configure(api_key=GEMINI_API_KEY)

try:
    model = genai.GenerativeModel("models/gemini-2.5-flash")
except:
    model = genai.GenerativeModel("models/gemini-flash-latest")

# Helper Functions
def extract_text(file):
    """Extracts text from TXT, PDF, DOCX, and PPTX files."""
    try:
        if file.name.endswith(".txt"):
            return file.read().decode("utf-8")
        elif file.name.endswith(".pdf"):
            reader = PdfReader(file)
            return " ".join([page.extract_text() for page in reader.pages if page.extract_text()])
        elif file.name.endswith(".docx"):
            doc = Document(file)
            return " ".join([p.text for p in doc.paragraphs])
        elif file.name.endswith(".pptx"):
            prs = Presentation(file)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return " ".join(text_runs)
        else:
            st.error("Unsupported file type.")
            return ""
    except Exception as e:
        st.error(f"Error reading file: {e}")
        return ""

def calculate_readability(text):
    """Calculates Flesch Reading Ease score."""
    return textstat.flesch_reading_ease(text)

def grammar_correction(text):
    st.info("â³ Sending text to Gemini for correction...")
    if len(text) > 30000:
        text = text[:30000]
        st.warning("Text truncated to 30000 chars for processing.")
        
    prompt = f"""
    You are a writing assistant. Correct grammatical errors and improve readability. Replace numbers with their corresponding words.
    Do not change the underlying meaning.
    Return ONLY the improved text.
    
    Text:
    {text}
    """
    try:
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        st.error(f"Gemini API Error: {e}")
        return text

def translate_text(text, target_language_code='en'):
    prompt = f"""
    Translate the following text into language code '{target_language_code}'.
    Ensure the translation is natural and suitable for being read aloud.
    Return ONLY the translated text.
    Text:
    {text}
    """
    try:
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        st.error(f"Translation Error: {e}")
        return text

def text_to_speech(text, lang="en", output_path="output_audio.mp3"):
    """Convert text to speech using gTTS."""
    try:
        clean_text = re.sub(r'\b\d+\b', lambda m: num2words(int(m.group()), lang=lang if lang in ['en', 'es', 'fr', 'de'] else 'en'), text)
    except:
        clean_text = text 
        
    try:
        tts = gTTS(text=clean_text, lang=lang)
        tts.save(output_path)
        return output_path
    except Exception as e:
        st.error(f"gTTS Error: {e}")
        return None

# def evaluate_tts_accuracy(original_text, audio_path, language_stt_code):
#     """Metrics Logic: WER calculation."""
#     r = sr.Recognizer()
#     wav_path = "temp_audio.wav"

#     try:
#         sound = AudioSegment.from_mp3(audio_path)
#         sound.export(wav_path, format="wav")

#         with sr.AudioFile(wav_path) as source:
#             audio_data = r.record(source)
#             transcribed_text = r.recognize_google(audio_data, language=language_stt_code)

#         wer = jiwer.wer(original_text, transcribed_text)
        
#         if os.path.exists(wav_path):
#             os.remove(wav_path)
            
#         return wer, transcribed_text

#     except Exception as e:
#         st.warning(f"Could not run TTS metrics: {e}")
#         return None, None

def normalize_text(text):
    """
    Removes punctuation, lowers case, and standardizes text 
    to make the comparison fair.
    """
    # Convert to lowercase
    text = text.lower()
    
    # Remove punctuation
    text = text.translate(str.maketrans('', '', string.punctuation))
    
    # Remove extra whitespace
    text = " ".join(text.split())
    
    return text

def evaluate_tts_accuracy(original_text, audio_path, language_stt_code):
    """
    Metrics Logic:
    1. Convert MP3 audio to WAV.
    2. Transcribe audio back to text (STT).
    3. Normalize both texts (remove punctuation/caps).
    4. Calculate WER.
    """
    r = sr.Recognizer()
    wav_path = "temp_audio.wav"

    try:
        # Convert MP3 to WAV
        sound = AudioSegment.from_mp3(audio_path)
        sound.export(wav_path, format="wav")

        # Transcribe
        with sr.AudioFile(wav_path) as source:
            audio_data = r.record(source)
            transcribed_text = r.recognize_google(audio_data, language=language_stt_code)

        # Normalize before comparing
        norm_original = normalize_text(original_text)
        norm_transcribed = normalize_text(transcribed_text)
        
        # Calculate WER on the CLEANED text
        wer = jiwer.wer(norm_original, norm_transcribed)
        
        # Cleanup
        if os.path.exists(wav_path):
            os.remove(wav_path)
            
        return wer, transcribed_text

    except Exception as e:
        st.warning(f"Could not run TTS metrics: {e}")
        return None, None

# Streamlit UI
st.set_page_config(page_title="AudioBook Gen & Metrics", layout="wide")

st.title("ðŸŽ§ Smart AudioBook Generator with QA Metrics")
st.markdown("Supported Formats: **.txt, .pdf, .docx, .pptx**")

with st.sidebar:
    st.header("Settings")
    selected_language = st.selectbox("Output Language", list(SUPPORTED_LANGUAGES.keys()))
    lang_config = SUPPORTED_LANGUAGES[selected_language]
    enable_metrics = st.checkbox("Enable QA Metrics (Slower)", value=True)

uploaded_file = st.file_uploader("Upload Document", type=["txt", "pdf", "docx", "pptx"])

if uploaded_file:
    with st.spinner("ðŸ“„ Extracting text..."):
        raw_text = extract_text(uploaded_file)
    
    if not raw_text.strip():
        st.error("File is empty or extraction failed.")
        st.stop()

    col1, col2 = st.columns(2)
    with col1:
        st.subheader("Raw Text")

        st.text_area("Preview", raw_text[:500] + "...", height=150, key="raw_preview")
        raw_score = calculate_readability(raw_text)
        st.caption(f"ðŸ“‰ Readability Score (Flesch): {raw_score:.2f}")

    if st.button("ðŸš€ Process & Generate Audio"):
        
        with st.spinner("âœ¨ Polishing grammar with Gemini..."):
            corrected_text = grammar_correction(raw_text)
        
        with col2:
            st.subheader("Corrected Text")

            st.text_area("Preview", corrected_text[:500] + "...", height=150, key="clean_preview")
            clean_score = calculate_readability(corrected_text)
            st.caption(f"ðŸ“ˆ Readability Score (Flesch): {clean_score:.2f}")
            
            delta = clean_score - raw_score
            if delta > 0:
                st.success(f"Readability improved by +{delta:.2f} points!")
            else:
                st.info("Readability remains similar.")

        target_code = lang_config["code"]
        stt_code = lang_config["stt_code"]
        
        final_text_for_audio = corrected_text
        
        # if target_code != "en":
        with st.spinner(f"ðŸŒ Translating to {selected_language}..."):
            final_text_for_audio = translate_text(corrected_text, target_code)
            st.text_area(f"Translated Text ({selected_language})", final_text_for_audio[:1000]+"...", height=150, key="trans_preview")

        audio_filename = f"audio_{target_code}.mp3"
        with st.spinner(f"ðŸ”Š Synthesizing Speech ({selected_language})..."):
            audio_path = text_to_speech(final_text_for_audio, lang=target_code, output_path=audio_filename)

        if audio_path:
            st.success("Audio Generation Complete!")
            st.audio(audio_path)
            
            with open(audio_path, "rb") as f:
                st.download_button("Download Audio", f, file_name=audio_filename)
            
            if enable_metrics:
                st.markdown("---")
                st.subheader("ðŸ“Š Quality Assurance Metrics")
                with st.spinner("Testing audio accuracy (Round-trip Transcribing)..."):
                    wer, transcribed = evaluate_tts_accuracy(final_text_for_audio, audio_path, stt_code)
                
                if wer is not None:
                    # Calculate accuracy
                    accuracy = max(0, (1 - wer) * 100)
                    
                    # Define detailed help strings
                    wer_help = """
                    Word Error Rate (WER) measures how different the audio is from the text.
                    a. 0.00 - 0.15: Excellent (Crystal clear).
                    b. 0.15 - 0.35: Good (Understandable, minor accent).
                    c. > 0.35: Low Quality (Too fast or unclear).
                    *Note: This is a strict computer test. 0.2 is often still very good for humans.*
                    """

                    acc_help = """
                    An estimate of how intelligible the audio is.
                    We convert the audio BACK to text to test it.
                    a. 85% - 100%: Perfect pronunciation.
                    b. 65% - 85%: Clear, but might have 'robot' accent quirks.
                    c. < 65%: Might be skipping words.
                    """

                    # Display Metrics
                    m_col1, m_col2, m_col3 = st.columns(3)
                    
                    m_col1.metric(
                        "Word Error Rate (WER)", 
                        f"{wer:.2f}", 
                        help=wer_help
                    )
                    
                    m_col2.metric(
                        "Est. Accuracy", 
                        f"{accuracy:.1f}%", 
                        help=acc_help
                    )
                    
                    # Visual Rating Logic
                    if wer <= 0.15:
                        m_col3.success("âœ… Excellent Quality")
                    elif wer <= 0.35:
                        m_col3.warning("âš ï¸ Acceptable Quality")
                    else:
                        m_col3.error("âŒ Low Intelligibility")

                    with st.expander("See what the computer heard (STT Transcript)"):
                        st.write(transcribed)

                        st.caption("Comparison between the text we sent to TTS vs. what the computer heard back.")
